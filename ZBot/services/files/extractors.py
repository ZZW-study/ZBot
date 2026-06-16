"""从 Office / PDF 字节中提取可发送给 LLM 的纯文本。

所有函数签名一致:`(raw_bytes) -> (text, error_or_None)`.
- text 非空:成功,作为文本块嵌入 message content。
- text 为空 + error 非空:失败,以 ``[文件名 解析失败: <reason>]`` 文本块告知模型,不抛 5xx。
"""
from __future__ import annotations

from io import BytesIO
from typing import Final

_MAX_PDF_PAGES: Final = 10
_MAX_TEXT_CHARS: Final = 120_000


def _truncate(text: str) -> str:
    if len(text) <= _MAX_TEXT_CHARS:
        return text
    return text[:_MAX_TEXT_CHARS] + "\n……(已截断)"


def extract_pdf(raw: bytes) -> tuple[str, str | None]:
    try:
        from pypdf import PdfReader
    except Exception as e:  # pragma: no cover
        return "", f"pypdf 不可用: {type(e).__name__}: {e}"
    try:
        reader = PdfReader(BytesIO(raw))
        n = min(len(reader.pages), _MAX_PDF_PAGES)
        text = "\n".join((reader.pages[i].extract_text() or "") for i in range(n))
        text = _truncate(text)
        if not text.strip():
            return "", "PDF 解析后无文本(可能为扫描件或加密文档)"
        return text, None
    except Exception as e:
        return "", f"PDF 解析失败: {type(e).__name__}: {e}"


def extract_docx(raw: bytes) -> tuple[str, str | None]:
    try:
        from docx import Document
    except Exception as e:  # pragma: no cover
        return "", f"python-docx 不可用: {type(e).__name__}: {e}"
    try:
        doc = Document(BytesIO(raw))
        text = "\n".join(p.text for p in doc.paragraphs)
        text = _truncate(text)
        if not text.strip():
            return "", "DOCX 解析后无文本"
        return text, None
    except Exception as e:
        return "", f"DOCX 解析失败: {type(e).__name__}: {e}"


def extract_xlsx(raw: bytes) -> tuple[str, str | None]:
    try:
        from openpyxl import load_workbook
    except Exception as e:  # pragma: no cover
        return "", f"openpyxl 不可用: {type(e).__name__}: {e}"
    try:
        wb = load_workbook(BytesIO(raw), data_only=True, read_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"## {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    parts.append(" | ".join(cells))
        text = _truncate("\n".join(parts))
        if not text.strip():
            return "", "XLSX 解析后无文本"
        return text, None
    except Exception as e:
        return "", f"XLSX 解析失败: {type(e).__name__}: {e}"


def extract_pptx(raw: bytes) -> tuple[str, str | None]:
    try:
        from pptx import Presentation
    except Exception as e:  # pragma: no cover
        return "", f"python-pptx 不可用: {type(e).__name__}: {e}"
    try:
        prs = Presentation(BytesIO(raw))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            parts.append(p.text)
        text = _truncate("\n".join(parts))
        if not text.strip():
            return "", "PPTX 解析后无文本"
        return text, None
    except Exception as e:
        return "", f"PPTX 解析失败: {type(e).__name__}: {e}"


# MIME -> 提取器 映射,给上层 agent_files 用
EXTRACTORS = {
    "application/pdf": extract_pdf,
    "application/msword": extract_docx,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
}
